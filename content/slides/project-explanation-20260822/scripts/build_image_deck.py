#!/usr/bin/env python3
"""画像版 PPTX を組み立てる（全スライドが 16:9 の画像 1 枚で構成される）。

スライドごとの画像は 3 系統:
  - 新規生成 `new-NN`: `gpt-image-2` の PNG（既定 `/tmp/claude/slide-images/`）を JPEG 化して取り込む
  - 実 UI    `shot-NN`: Playwright で撮った PNG（`images/` に既にある）
  - 既存流用        : `docs/infographics/*.webp` を JPEG 化して取り込む

PNG のままだと 1 枚約 2MB で `docs/infographics/README.md` の方針（PNG はリポジトリに入れない）に
反するため、リポジトリに置くのは JPEG（品質は下記 QUALITY）に統一する。`python-pptx` は WebP を埋め込めないので、
既存インフォグラフィックも JPEG へ変換してから貼る。

品質は 78。88 だと 19 枚を束ねた PPTX が 5.5MB になり、`tools/self_review_check.py` の巨大ファイル
検査（5.0MB 超で PR 作成をブロック）に掛かる。平坦な色面のイラストなので、78 でも投影時の劣化は
判別できない。
"""

import json
import re
import shutil
import sys
from pathlib import Path

from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
REPO = Path(__file__).resolve().parents[4]
sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

PLAN = ROOT / "content" / "slides_plan.json"
IMAGES = ROOT / "images"
STAGE = Path("/tmp/claude/slide-deck-stage")
RAW_NEW = Path("/tmp/claude/slide-images")
OUT = ROOT / "output" / "gem-hunter.pptx"

SIZE = (1536, 864)  # 完全な 16:9（既存インフォグラフィックと同じ）
QUALITY = 78        # PPTX 全体を 5.0MB 未満に収めるための圧縮率（上のドキストリング参照）


def to_jpeg(src: Path, dest: Path) -> Path:
    """任意の画像を 1536x864 の JPEG（quality 88）へ正規化する。"""
    with Image.open(src) as im:
        im = im.convert("RGB")
        if im.size != SIZE:
            im = im.resize(SIZE, Image.LANCZOS)
        dest.parent.mkdir(parents=True, exist_ok=True)
        im.save(dest, "JPEG", quality=QUALITY, optimize=True)
    return dest


def resolve(slide: dict) -> Path:
    """スライドの `visual` 記述から、貼り付ける画像の実体パスを決める。"""
    visual = slide["visual"]
    new_id = re.search(r"(new-\d+)", visual)
    if new_id:
        raw = RAW_NEW / f"{new_id.group(1)}.png"
        if not raw.exists():
            raise FileNotFoundError(f"生成画像がない: {raw}")
        return to_jpeg(raw, IMAGES / f"{new_id.group(1)}.jpg")
    shot_id = re.search(r"(shot-\d+)", visual)
    if shot_id:
        found = sorted(IMAGES.glob(f"{shot_id.group(1)}-*.png"))
        if not found:
            raise FileNotFoundError(f"スクリーンショットがない: {shot_id.group(1)}")
        return to_jpeg(found[0], STAGE / f"{shot_id.group(1)}.jpg")
    reused = re.search(r"(docs/infographics/[\w.-]+\.webp)", visual)
    if reused:
        src = REPO / reused.group(1)
        if not src.exists():
            raise FileNotFoundError(f"既存インフォグラフィックがない: {src}")
        return to_jpeg(src, STAGE / (src.stem + ".jpg"))
    raise ValueError(f'スライド {slide["no"]} の visual を解決できない: {visual}')


def main() -> int:
    plan = json.loads(PLAN.read_text(encoding="utf-8"))
    shutil.rmtree(STAGE, ignore_errors=True)
    prs = Presentation()
    prs.slide_width = Emu(int(12192000))  # 13.333in = 16:9
    prs.slide_height = Emu(int(6858000))  # 7.5in
    for slide in plan["slides"]:
        path = resolve(slide)
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_picture(
            str(path), left=Emu(0), top=Emu(0),
            width=prs.slide_width, height=prs.slide_height,
        )
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Image PPTX generated: {OUT} ({len(plan['slides'])} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
