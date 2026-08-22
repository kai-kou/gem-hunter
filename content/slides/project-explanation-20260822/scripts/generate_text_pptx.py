#!/usr/bin/env python3
"""slides_plan.json からテキスト版 PPTX を生成する。

レイアウト関数は `pptx_template.py`（参照リポジトリのテンプレートをそのまま取り込んだもの）を
import して使い、配色だけを既存インフォグラフィック 13 枚（ネイビー / ティール / コーラル /
生成りの紙地）に合わせて差し替える。画像版と並べたときにトーンが割れないようにするため。
"""

import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(Path(__file__).resolve().parent))

import pptx_template as tpl  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR  # noqa: E402
from pptx.util import Pt  # noqa: E402

PLAN = ROOT / "content" / "slides_plan.json"
OUT = ROOT / "output" / "gem-hunter_text.pptx"

# 選択テーマ: 既存インフォグラフィックと同じ「生成りの紙地 + ネイビー + ティール」。
# 理由: 対象読者が開発者で、同一デッキ内に既存グラレコ画像を混在させるため配色を揃える。
tpl.C_BG = RGBColor(0xFA, 0xF6, 0xEC)
tpl.C_HEADER = RGBColor(0x14, 0x36, 0x4F)
tpl.C_ACCENT = RGBColor(0x14, 0x74, 0x6F)
tpl.C_TEXT = RGBColor(0x1B, 0x2A, 0x3A)
tpl.C_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

SUMMARY_CARD_TITLES = ["実利用を起点に", "残差で測る", "層を足す条件", "殺せる記録"]


def _auto_font_size(item_count: int) -> float:
    """本文のフォントサイズ。

    テンプレート既定（6 点で 12pt）は、1 点 40〜60 字という本デッキの分量だと
    16:9 の下半分が空いたまま小さく詰まって見える。投影時の可読性を優先して底上げする。
    """
    if item_count <= 3:
        return 18.0
    if item_count <= 5:
        return 16.0
    if item_count <= 6:
        return 15.0
    return 14.0


tpl._auto_font_size = _auto_font_size


def to_slide_data(slide: dict) -> dict:
    """スライド定義（slides_plan.json）を、テンプレートが解釈する形へ変換する。

    レイアウトの指定は `layout` フィールドが正本（スライド番号にハードコードしない）。
    """
    layout = slide.get("layout", "bullets")
    if layout == "title":
        return {
            "type": "title",
            "title": slide["elements"][0],
            "subtitle": " / ".join(slide["elements"][1:]),
        }
    if layout == "summary":
        # カードは SUMMARY_CARD_TITLES の数だけしか作れない。あふれた要素を黙って捨てると
        # 構成にある文言が PPTX から消えるので、ここで明示的に失敗させる。
        if len(slide["elements"]) > len(SUMMARY_CARD_TITLES):
            raise ValueError(
                f'スライド {slide["no"]}: summary レイアウトは {len(SUMMARY_CARD_TITLES)} 点までしか'
                f'描けないが本文が {len(slide["elements"])} 点ある。bullets に変えるか本文を減らすこと。'
            )
        cards = [
            {"title": t, "body": b}
            for t, b in zip(SUMMARY_CARD_TITLES, slide["elements"])
        ]
        return {"type": "summary", "header": slide["title"], "cards": cards}
    if layout != "bullets":
        # 未対応の値をサイレントに箇条書きへ落とすと、意図したレイアウトが失われたことに気づけない。
        raise ValueError(f'スライド {slide["no"]}: 未対応の layout: {layout}')
    return {"type": "bullets", "header": slide["title"], "points": slide["elements"]}


def space_out_bullets(path: Path) -> None:
    """本文テキストボックスの行間を空け、ブロックを上下中央に寄せる。

    テンプレート既定は行間ゼロ・上寄せで、16:9 の下半分が空いたまま上に詰まる。
    段落間に余白を入れ、テキストボックスの垂直位置を中央にして紙面を使い切る。
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX or not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if len(tf.paragraphs) < 2:
                continue
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in tf.paragraphs:
                # テンプレート既定の行頭記号「▸」は日本語フォントによっては字形を持たず
                # 豆腐（□）になる。どの環境でも出る「・」に置き換える。
                for run in para.runs:
                    if run.text.startswith("▸"):
                        run.text = "・" + run.text.lstrip("▸ ")
                size = para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else 14
                para.space_after = Pt(size * 0.75)
                para.line_spacing = 1.25
    prs.save(str(path))


def main() -> int:
    plan = json.loads(PLAN.read_text(encoding="utf-8"))
    slides_data = [to_slide_data(s) for s in plan["slides"]]
    OUT.parent.mkdir(parents=True, exist_ok=True)
    tpl.create_text_pptx(slides_data, str(OUT))
    space_out_bullets(OUT)
    print(f"spacing applied: {OUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
