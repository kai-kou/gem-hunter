"""テキスト版 PPTX のレイアウト関数群。

出典: kai-kou/qiita-bash-lt-2026 の `.claude/skills/slides/references/PPTX_TEMPLATE.md`
「テキスト版PPTXスクリプトテンプレート」（取得日 2026-08-22・MIT License）。
テンプレートを本文へ書き写さず、コード片としてそのまま取り込んで import して使う
（参照リポジトリ `.claude/rules/lessons.md` の「テンプレートは丸コピーせず import する」に従う）。

カラー定数（C_BG / C_HEADER / C_ACCENT / C_TEXT / C_CARD_BG）は呼び出し側が
モジュール属性として差し替える前提。既定値は参照テンプレートのまま。
"""
import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ===== カラー定数（コンテンツ分析で選択したテーマを設定）=====
C_BG        = RGBColor(0xF0, 0xF4, 0xF8)   # 背景色
C_HEADER    = RGBColor(0x1C, 0x35, 0x57)   # ヘッダーバー
C_ACCENT    = RGBColor(0x25, 0x63, 0xEB)   # アクセント
C_TEXT      = RGBColor(0x1E, 0x29, 0x3B)   # 本文テキスト
C_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)  # ヘッダー上テキスト（白）
C_CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)   # カード背景


# ===== 単位変換ユーティリティ =====
def mm(value: float) -> int:
    """mm → EMU（1mm = 36000 EMU）"""
    return int(value * 36000)

def cm(value: float) -> int:
    return int(value * 360000)

# スライドサイズ定数（16:9）
SLIDE_W = mm(337.0)
SLIDE_H = mm(190.0)

# 標準レイアウト定数
HEADER_H    = mm(28)
MARGIN_L    = mm(17)
MARGIN_R    = mm(17)
MARGIN_T    = mm(8)   # ヘッダー内テキスト上余白
CONTENT_TOP = HEADER_H + mm(10)
CONTENT_H   = SLIDE_H - HEADER_H - mm(15)
CONTENT_W   = SLIDE_W - MARGIN_L - MARGIN_R


# ===== ユーティリティ関数 =====

def set_bg(slide, color: RGBColor):
    """スライド背景色を設定する"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text: str):
    """ヘッダーバーとタイトルテキストを追加する"""
    # ヘッダーバー（矩形）
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left=0, top=0,
        width=SLIDE_W, height=HEADER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_HEADER
    bar.line.fill.background()  # 枠線なし

    # ヘッダーテキスト
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.color.rgb = C_TEXT_LIGHT
    run.font.size = Pt(18)
    run.font.bold = True
    _set_jp_font(run)


def add_shadow(shape):
    """白背景カードに軽い影を付ける"""
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '101600')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '8100000')
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', '000000')
    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', '10000')


def _set_jp_font(run, size_pt: float = None, bold: bool = None):
    """日本語対応フォントをセットする（游ゴシック + Calibri）"""
    run.font.name = 'Calibri'
    if size_pt:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', '游ゴシック')


def add_text_box(slide, left, top, width, height,
                 text, size_pt=14, color=None, bold=False,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    """テキストボックスを追加する汎用関数"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.color.rgb = color or C_TEXT
    _set_jp_font(run, size_pt, bold)
    return txBox


def add_bullet_paragraph(tf, text, size_pt=13, level=0, color=None):
    """箇条書きの段落をテキストフレームに追加する"""
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.color.rgb = color or C_TEXT
    _set_jp_font(run, size_pt)
    return p


def add_card(slide, left, top, width, height, with_shadow=True):
    """白背景カードを追加する（カード内にテキストを重ねて配置する）"""
    card = slide.shapes.add_shape(
        1, left=left, top=top, width=width, height=height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(0.5)
    if with_shadow:
        add_shadow(card)
    return card


# ===== レイアウト別スライド生成関数 =====

def make_title_slide(prs, title: str, subtitle: str):
    """タイトルスライド: 中央揃え大タイトル + サブタイトル"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # アクセントバー（上部）
    bar = slide.shapes.add_shape(1, left=0, top=0, width=SLIDE_W, height=mm(8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    # タイトル
    add_text_box(
        slide,
        left=mm(30), top=mm(50),
        width=SLIDE_W - mm(60), height=mm(60),
        text=title, size_pt=40, color=C_HEADER,
        bold=True, align=PP_ALIGN.CENTER
    )
    # サブタイトル
    add_text_box(
        slide,
        left=mm(30), top=mm(118),
        width=SLIDE_W - mm(60), height=mm(25),
        text=subtitle, size_pt=18, color=RGBColor(0x64, 0x74, 0x8B),
        align=PP_ALIGN.CENTER
    )
    return slide


def make_bullets_slide(prs, header: str, points: list[str]):
    """箇条書きスライド: ヘッダーバー + 本文箇条書き"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    # コンテンツカード
    card = add_card(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, CONTENT_H)
    txBox = slide.shapes.add_textbox(
        MARGIN_L + mm(8), CONTENT_TOP + mm(8),
        CONTENT_W - mm(16), CONTENT_H - mm(16)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"▸  {point}"
        run.font.color.rgb = C_TEXT
        _set_jp_font(run, _auto_font_size(len(points)))
        p.space_after = Pt(4)
    return slide


def make_two_column_slide(prs, header: str,
                           left_title: str, left_points: list[str],
                           right_title: str, right_points: list[str]):
    """2カラムスライド: ヘッダーバー + 左右カラム"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    col_w = (CONTENT_W - mm(6)) // 2
    gutter = mm(6)

    for i, (col_title, col_points) in enumerate(
        [(left_title, left_points), (right_title, right_points)]
    ):
        col_left = MARGIN_L + i * (col_w + gutter)
        # カラムカード
        add_card(slide, col_left, CONTENT_TOP, col_w, CONTENT_H)
        # カラムタイトル
        title_bg = slide.shapes.add_shape(
            1, left=col_left, top=CONTENT_TOP,
            width=col_w, height=mm(14)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = C_ACCENT
        title_bg.line.fill.background()
        add_text_box(
            slide, col_left + mm(4), CONTENT_TOP + mm(1),
            col_w - mm(8), mm(12),
            col_title, size_pt=13, color=C_TEXT_LIGHT, bold=True
        )
        # 本文
        body_top = CONTENT_TOP + mm(18)
        body_h = CONTENT_H - mm(22)
        txBox = slide.shapes.add_textbox(
            col_left + mm(6), body_top, col_w - mm(12), body_h
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, pt in enumerate(col_points):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {pt}"
            run.font.color.rgb = C_TEXT
            _set_jp_font(run, _auto_font_size(len(col_points)))
            p.space_after = Pt(3)
    return slide


def make_key_message_slide(prs, header: str,
                            key_message: str, sub_points: list[str]):
    """キーメッセージ強調スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    # キーメッセージ（大きく中央）
    msg_h = mm(50)
    add_card(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, msg_h)
    add_text_box(
        slide, MARGIN_L + mm(8), CONTENT_TOP + mm(8),
        CONTENT_W - mm(16), msg_h - mm(16),
        key_message, size_pt=24, color=C_HEADER,
        bold=True, align=PP_ALIGN.CENTER
    )

    # サブポイント
    sub_top = CONTENT_TOP + msg_h + mm(8)
    sub_h = CONTENT_H - msg_h - mm(8)
    txBox = slide.shapes.add_textbox(MARGIN_L, sub_top, CONTENT_W, sub_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(sub_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"▸  {pt}"
        run.font.color.rgb = C_TEXT
        _set_jp_font(run, 13)
        p.space_after = Pt(4)
    return slide


def make_summary_slide(prs, header: str, cards: list[dict]):
    """まとめスライド: アイコン付き要点カード（最大4枚）
    cards: [{"title": str, "body": str}, ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    n = len(cards)
    card_w = (CONTENT_W - mm(5) * (n - 1)) // n
    for i, card_data in enumerate(cards):
        cx = MARGIN_L + i * (card_w + mm(5))
        add_card(slide, cx, CONTENT_TOP, card_w, CONTENT_H)
        # カードタイトル
        t_bg = slide.shapes.add_shape(
            1, left=cx, top=CONTENT_TOP, width=card_w, height=mm(16)
        )
        t_bg.fill.solid()
        t_bg.fill.fore_color.rgb = C_ACCENT
        t_bg.line.fill.background()
        add_text_box(
            slide, cx + mm(3), CONTENT_TOP + mm(2),
            card_w - mm(6), mm(12),
            card_data["title"], size_pt=12, color=C_TEXT_LIGHT,
            bold=True, align=PP_ALIGN.CENTER
        )
        # カード本文
        add_text_box(
            slide, cx + mm(5), CONTENT_TOP + mm(20),
            card_w - mm(10), CONTENT_H - mm(26),
            card_data["body"], size_pt=12, color=C_TEXT
        )
    return slide


def make_section_slide(prs, section_title: str, description: str = ""):
    """セクション区切りスライド: 左端アクセントバー + 中央揃えタイトル + 1行説明"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # 左端アクセントバー
    bar = slide.shapes.add_shape(1, left=0, top=0, width=mm(8), height=SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    # セクションタイトル
    add_text_box(
        slide,
        left=mm(30), top=mm(68),
        width=SLIDE_W - mm(60), height=mm(55),
        text=section_title, size_pt=32, color=C_HEADER,
        bold=True, align=PP_ALIGN.CENTER
    )
    if description:
        add_text_box(
            slide,
            left=mm(30), top=mm(130),
            width=SLIDE_W - mm(60), height=mm(22),
            text=description, size_pt=16,
            color=RGBColor(0x64, 0x74, 0x8B),
            align=PP_ALIGN.CENTER
        )
    return slide


def make_three_column_slide(prs, header: str,
                             col1_title: str, col1_points: list[str],
                             col2_title: str, col2_points: list[str],
                             col3_title: str, col3_points: list[str]):
    """3カラムスライド: ヘッダーバー + 3等分カラム（各2〜3項目）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    gutter = mm(5)
    col_w = (CONTENT_W - gutter * 2) // 3

    for i, (col_title, col_points) in enumerate([
        (col1_title, col1_points),
        (col2_title, col2_points),
        (col3_title, col3_points),
    ]):
        col_left = MARGIN_L + i * (col_w + gutter)
        add_card(slide, col_left, CONTENT_TOP, col_w, CONTENT_H)

        # カラムタイトルバー
        t_bg = slide.shapes.add_shape(
            1, left=col_left, top=CONTENT_TOP, width=col_w, height=mm(14)
        )
        t_bg.fill.solid()
        t_bg.fill.fore_color.rgb = C_ACCENT
        t_bg.line.fill.background()
        add_text_box(
            slide, col_left + mm(3), CONTENT_TOP + mm(1),
            col_w - mm(6), mm(12),
            col_title, size_pt=12, color=C_TEXT_LIGHT,
            bold=True, align=PP_ALIGN.CENTER
        )

        # 本文
        body_top = CONTENT_TOP + mm(18)
        body_h = CONTENT_H - mm(22)
        txBox = slide.shapes.add_textbox(
            col_left + mm(5), body_top, col_w - mm(10), body_h
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, pt in enumerate(col_points):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {pt}"
            run.font.color.rgb = C_TEXT
            _set_jp_font(run, _auto_font_size(len(col_points) + 1))
            p.space_after = Pt(3)
    return slide


def make_comparison_table_slide(prs, header: str,
                                  columns: list[str],
                                  rows: list[dict]):
    """比較表スライド: ヘッダーバー + テーブル（最大5列）
    columns: 比較軸ラベルのリスト（例: ["手動", "Claude"]）
    rows: [{"label": "項目名", "values": ["手動の値", "Claudeの値"]}, ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    n_rows = len(rows)
    row_h = min(mm(22), (CONTENT_H - mm(4)) // (n_rows + 1))
    col_w_label = mm(55)
    col_w_val = (CONTENT_W - col_w_label) // len(columns)

    # ヘッダー行
    header_items = [""] + columns
    for i, col_name in enumerate(header_items):
        if i == 0:
            cx, cw = MARGIN_L, col_w_label
        else:
            cx = MARGIN_L + col_w_label + col_w_val * (i - 1)
            cw = col_w_val
        hdr = slide.shapes.add_shape(
            1, left=cx, top=CONTENT_TOP, width=cw, height=row_h
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = C_HEADER if i > 0 else RGBColor(0xCB, 0xD5, 0xE1)
        hdr.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        hdr.line.width = Pt(0.5)
        if col_name:
            add_text_box(
                slide, cx + mm(2), CONTENT_TOP + mm(2), cw - mm(4), row_h - mm(4),
                col_name, size_pt=12, color=C_TEXT_LIGHT,
                bold=True, align=PP_ALIGN.CENTER
            )

    # データ行
    for r, row in enumerate(rows):
        row_top = CONTENT_TOP + row_h * (r + 1)
        row_bg = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFC)

        # ラベル列
        lbl = slide.shapes.add_shape(
            1, left=MARGIN_L, top=row_top, width=col_w_label, height=row_h
        )
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        lbl.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        lbl.line.width = Pt(0.5)
        add_text_box(
            slide, MARGIN_L + mm(3), row_top + mm(2),
            col_w_label - mm(6), row_h - mm(4),
            row["label"], size_pt=12, color=C_TEXT, bold=True
        )

        # 値列
        for c, val in enumerate(row.get("values", [])):
            cell_left = MARGIN_L + col_w_label + col_w_val * c
            cell = slide.shapes.add_shape(
                1, left=cell_left, top=row_top, width=col_w_val, height=row_h
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            cell.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
            cell.line.width = Pt(0.5)
            add_text_box(
                slide, cell_left + mm(3), row_top + mm(2),
                col_w_val - mm(6), row_h - mm(4),
                val, size_pt=11, color=C_TEXT, align=PP_ALIGN.CENTER
            )
    return slide


def make_timeline_slide(prs, header: str, milestones: list[dict]):
    """タイムラインスライド: 横並びマイルストーン（最大6個）
    milestones: [{"label": "Step 1", "title": "要件定義", "desc": "概要1行"}, ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_header_bar(slide, header)

    n = len(milestones)
    item_w = CONTENT_W // n
    center_y = CONTENT_TOP + CONTENT_H // 2

    # 横ライン
    line = slide.shapes.add_shape(
        1,
        left=MARGIN_L, top=center_y - mm(1),
        width=CONTENT_W, height=mm(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

    for i, ms in enumerate(milestones):
        cx = MARGIN_L + item_w * i + item_w // 2

        # 丸アイコン
        dot_r = mm(7)
        dot = slide.shapes.add_shape(
            9,  # OVAL
            left=cx - dot_r, top=center_y - dot_r,
            width=dot_r * 2, height=dot_r * 2
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()

        # ラベル（上）
        if ms.get("label"):
            add_text_box(
                slide,
                cx - item_w // 2 + mm(3), center_y - mm(30),
                item_w - mm(6), mm(18),
                ms["label"], size_pt=10,
                color=RGBColor(0x94, 0xA3, 0xB8),
                align=PP_ALIGN.CENTER
            )

        # タイトル（下）
        add_text_box(
            slide,
            cx - item_w // 2 + mm(3), center_y + mm(12),
            item_w - mm(6), mm(20),
            ms["title"], size_pt=13, color=C_TEXT,
            bold=True, align=PP_ALIGN.CENTER
        )

        # 説明（タイトル下）
        if ms.get("desc"):
            add_text_box(
                slide,
                cx - item_w // 2 + mm(3), center_y + mm(33),
                item_w - mm(6), mm(22),
                ms["desc"], size_pt=10,
                color=RGBColor(0x64, 0x74, 0x8B),
                align=PP_ALIGN.CENTER
            )
    return slide


def _auto_font_size(item_count: int) -> float:
    """項目数に応じてフォントサイズを自動調整"""
    if item_count <= 3:
        return 16.0
    elif item_count <= 5:
        return 14.0
    elif item_count <= 7:
        return 12.0
    else:
        return 11.0


# ===== メイン: スライドデータを受け取ってPPTXを生成する =====

def create_text_pptx(slides_data: list[dict], output_path: str):
    """
    slides_data: [
        {"type": "title",            "title": "...", "subtitle": "..."},
        {"type": "section",          "section_title": "...", "description": "（省略可）"},
        {"type": "bullets",          "header": "...", "points": [...]},
        {"type": "two-column",       "header": "...",
         "left_title": "...", "left_points": [...],
         "right_title": "...", "right_points": [...]},
        {"type": "three-column",     "header": "...",
         "col1_title": "...", "col1_points": [...],
         "col2_title": "...", "col2_points": [...],
         "col3_title": "...", "col3_points": [...]},
        {"type": "comparison-table", "header": "...",
         "columns": ["選択肢A", "選択肢B"],
         "rows": [{"label": "項目1", "values": ["A値", "B値"]}, ...]},
        {"type": "key-message",      "header": "...", "key_message": "...", "sub_points": [...]},
        {"type": "timeline",         "header": "...",
         "milestones": [{"label": "Step 1", "title": "要件定義", "desc": "..."}, ...]},
        {"type": "summary",          "header": "...", "cards": [{"title": "...", "body": "..."}, ...]},
    ]
    """
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    layout_fn = {
        "title":            lambda s: make_title_slide(prs, s["title"], s["subtitle"]),
        "section":          lambda s: make_section_slide(
            prs, s["section_title"], s.get("description", "")
        ),
        "bullets":          lambda s: make_bullets_slide(prs, s["header"], s["points"]),
        "two-column":       lambda s: make_two_column_slide(
            prs, s["header"],
            s["left_title"], s["left_points"],
            s["right_title"], s["right_points"]
        ),
        "three-column":     lambda s: make_three_column_slide(
            prs, s["header"],
            s["col1_title"], s["col1_points"],
            s["col2_title"], s["col2_points"],
            s["col3_title"], s["col3_points"]
        ),
        "comparison-table": lambda s: make_comparison_table_slide(
            prs, s["header"], s["columns"], s["rows"]
        ),
        "key-message":      lambda s: make_key_message_slide(
            prs, s["header"], s["key_message"], s["sub_points"]
        ),
        "timeline":         lambda s: make_timeline_slide(prs, s["header"], s["milestones"]),
        "summary":          lambda s: make_summary_slide(prs, s["header"], s["cards"]),
    }

    for slide_data in slides_data:
        fn = layout_fn.get(slide_data["type"])
        if fn:
            fn(slide_data)
        else:
            # フォールバック: 箇条書き
            make_bullets_slide(prs, slide_data.get("header", ""), slide_data.get("points", []))

    prs.save(output_path)
    print(f"Text PPTX generated: {output_path} ({len(slides_data)} slides)")


if __name__ == "__main__":
    config = json.loads(sys.argv[1])
    create_text_pptx(
        slides_data=config["slides"],
        output_path=config["output_path"],
    )