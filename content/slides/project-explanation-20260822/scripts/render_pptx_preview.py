#!/usr/bin/env python3
"""生成した PPTX を PNG に描き起こしてセルフレビュー・ユーザーレビューに使う。

この環境の LibreOffice は `libreoffice-core` だけで文書フィルタ（Impress / Writer）が入っておらず、
`soffice --convert-to pdf` は .txt すら「source file could not be loaded」で失敗する。
そこで PPTX の中身（矩形とテキストボックス）を python-pptx で読み、Pillow で描き直す。
テンプレートが使う図形は矩形とテキストのみなので、これで実物と等価な絵になる。
"""

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

FONT = "/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf"
OUT_W = 1536


def load_font(size_px: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT, max(size_px, 8))


def rgb(color, default=(30, 30, 30)):
    try:
        if color is None or color.type is None:
            return default
        return tuple(color.rgb)
    except (AttributeError, TypeError, ValueError):
        return default


def draw_shape(drw: ImageDraw.ImageDraw, shape, k: float) -> None:
    x0, y0 = shape.left * k, shape.top * k
    x1, y1 = (shape.left + shape.width) * k, (shape.top + shape.height) * k
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # MSO_FILL.SOLID
            drw.rectangle([x0, y0, x1, y1], fill=rgb(fill.fore_color, (255, 255, 255)))
    except (AttributeError, TypeError, ValueError):
        pass
    if not shape.has_text_frame:
        return
    # 実物の見え方に合わせる: 図形（矩形）の文字は上下中央、テキストボックスは
    # vertical_anchor の指定に従う（既定は上寄せ）。
    lines_all = _layout_lines(drw, shape, k, x1 - x0 - 8)
    total_h = sum(h for _, _, h in lines_all)
    anchor_middle = shape.shape_type != 17 or str(shape.text_frame.vertical_anchor) == "MIDDLE (3)"
    y = y0 + max(4.0, ((y1 - y0) - total_h) / 2) if anchor_middle else y0 + 4
    centered = any(str(p.alignment) == "CENTER (2)" for p in shape.text_frame.paragraphs)
    for text, font, h in lines_all:
        x = x0 + 4
        if centered:
            x = x0 + ((x1 - x0) - drw.textlength(text, font=font)) / 2
        drw.text((x, y), text, font=font, fill=_line_color(shape, text))
        y += h
    return


def _line_color(shape, _text):
    for para in shape.text_frame.paragraphs:
        if para.runs:
            return rgb(para.runs[0].font.color, (30, 30, 30))
    return (30, 30, 30)


def _layout_lines(drw, shape, k, max_w):
    """段落を折り返し、(テキスト, フォント, 行高) のリストにする。"""
    out = []
    for para in shape.text_frame.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text:
            out.append(("", load_font(12), 10))
            continue
        run = para.runs[0]
        size_pt = run.font.size.pt if run.font.size else 12
        font = load_font(int(size_pt * k * 12700))
        gap = para.space_after.pt * k * 12700 if para.space_after else 0
        line = ""
        buf = []
        for ch in text:
            if drw.textlength(line + ch, font=font) > max_w and line:
                buf.append(line)
                line = ch
            else:
                line += ch
        buf.append(line)
        for i, ln in enumerate(buf):
            h = font.size * 1.35 + (gap if i == len(buf) - 1 else 0)
            out.append((ln, font, h))
    return out


def _unused(drw, shape, k, x0, y0, x1, y1):
    y = y0 + 4
    for para in shape.text_frame.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text:
            y += 8
            continue
        run = para.runs[0]
        size_pt = run.font.size.pt if run.font.size else 12
        color = rgb(run.font.color, (30, 30, 30))
        font = load_font(int(size_pt * k * 12700))
        # 折り返し（描画幅に収まる文字数で機械的に割る）
        max_w = x1 - x0 - 8
        line, lines = "", []
        for ch in text:
            if drw.textlength(line + ch, font=font) > max_w and line:
                lines.append(line)
                line = ch
            else:
                line += ch
        lines.append(line)
        for ln in lines:
            drw.text((x0 + 4, y), ln, font=font, fill=color)
            y += font.size * 1.35


def render(pptx_path: Path, out_dir: Path) -> int:
    prs = Presentation(str(pptx_path))
    k = OUT_W / prs.slide_width          # EMU → px の倍率
    out_h = int(prs.slide_height * k)
    out_dir.mkdir(parents=True, exist_ok=True)
    for i, slide in enumerate(prs.slides, 1):
        bg = (255, 255, 255)
        try:
            fill = slide.background.fill
            if fill.type == 1:
                bg = tuple(fill.fore_color.rgb)
        except (AttributeError, TypeError, ValueError):
            pass
        img = Image.new("RGB", (OUT_W, out_h), bg)
        drw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            draw_shape(drw, shape, k)
        path = out_dir / f"text-{i:02d}.png"
        img.save(path)
    print(f"rendered {len(prs.slides._sldIdLst)} slides -> {out_dir}")
    return 0


if __name__ == "__main__":
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2])
    sys.exit(render(src, dst))
