#!/usr/bin/env python3
"""テキスト版 PPTX に構成の文言が漏れなく入っているかを機械照合する（Step 4 セルフレビュー）。

参照ワークフローの Step 4「内容整合チェック: 構成マークダウンの全テキストが PPTX に
正しく反映されているか照合する」を、目視ではなく機械で行う。
"""

import json
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PLAN = ROOT / "content" / "slides_plan.json"
PPTX = ROOT / "output" / "gem-hunter_text.pptx"


def norm(text: str) -> str:
    """行頭記号・空白の違いを無視して比較する。"""
    text = unicodedata.normalize("NFKC", text)
    return "".join(text.split()).lstrip("・•▸-")


def slide_text(slide) -> str:
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            out.append(shape.text_frame.text)
    return norm("\n".join(out))


def main() -> int:
    plan = json.loads(PLAN.read_text(encoding="utf-8"))
    prs = Presentation(str(PPTX))
    slides = list(prs.slides)
    missing = []
    if len(slides) != len(plan["slides"]):
        print(f"NG: 枚数が違う（構成 {len(plan['slides'])} / PPTX {len(slides)}）", file=sys.stderr)
        return 1
    for spec, slide in zip(plan["slides"], slides):
        body = slide_text(slide)
        wanted = list(spec["elements"])
        if spec.get("layout") != "title":
            wanted.append(spec["title"])
        for text in wanted:
            if norm(text) not in body:
                missing.append((spec["no"], text))
    if missing:
        for no, text in missing:
            print(f"NG: スライド {no} に反映されていない: {text}", file=sys.stderr)
        return 1
    print(f"OK: {len(slides)} 枚すべての見出し・本文が PPTX に反映されている")
    return 0


if __name__ == "__main__":
    sys.exit(main())
